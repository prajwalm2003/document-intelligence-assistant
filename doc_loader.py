# document_loader.py
# WHY: Instead of cramming all loaders into rag_pipeline.py,
# we separate concerns — clean software design (SE teachers love this)
# Each file type has its own loader function

from langchain_community.document_loaders import PyPDFLoader
from langchain_core.documents import Document
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
import os

def load_pdf(file_path):
    """Load PDF using LangChain's PyPDFLoader"""
    loader = PyPDFLoader(file_path)
    docs = loader.load()
    for doc in docs:
        doc.metadata["file_type"] = "PDF"
    return docs

def load_docx(file_path):
    """
    Load Word documents
    WHY: python-docx reads .docx files paragraph by paragraph
    We combine all paragraphs into one Document object per file
    """
    docx = DocxDocument(file_path)
    full_text = ""
    for para in docx.paragraphs:
        if para.text.strip():  # Skip empty paragraphs
            full_text += para.text + "\n"

    # Also extract tables from Word docs
    for table in docx.tables:
        for row in table.rows:
            row_text = " | ".join([cell.text.strip() for cell in row.cells])
            full_text += row_text + "\n"

    return [Document(
        page_content=full_text,
        metadata={"source": file_path, "file_type": "DOCX", "page": 0,
                  "source_file": os.path.basename(file_path)}
    )]

def load_txt(file_path):
    """
    Load plain text files
    WHY: Simplest loader — just read the file
    Good for logs, notes, code files
    """
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return [Document(
        page_content=text,
        metadata={"source": file_path, "file_type": "TXT", "page": 0,
                  "source_file": os.path.basename(file_path)}
    )]

def load_csv_excel(file_path):
    """
    Load CSV and Excel files
    WHY: pandas reads tabular data — we convert each row to text
    so Gemini can understand and answer questions about the data
    """
    if file_path.endswith(".csv"):
        df = pd.read_csv(file_path)
    else:
        df = pd.read_excel(file_path)

    # Convert dataframe to readable text chunks
    # WHY: LLMs understand text better than raw tables
    docs = []
    
    # First chunk — column overview
    col_info = f"This document has {len(df)} rows and {len(df.columns)} columns.\n"
    col_info += f"Columns are: {', '.join(df.columns.tolist())}\n\n"
    col_info += f"First few rows:\n{df.head(5).to_string()}"
    
    docs.append(Document(
        page_content=col_info,
        metadata={"source": file_path, "file_type": "CSV/Excel",
                  "page": 0, "source_file": os.path.basename(file_path)}
    ))

    # Split remaining rows into chunks of 50
    chunk_size = 50
    for i in range(0, len(df), chunk_size):
        chunk = df.iloc[i:i+chunk_size]
        text = f"Rows {i} to {i+len(chunk)}:\n{chunk.to_string()}"
        docs.append(Document(
            page_content=text,
            metadata={"source": file_path, "file_type": "CSV/Excel",
                      "page": i//chunk_size, "source_file": os.path.basename(file_path)}
        ))
    return docs

def load_pptx(file_path):
    """
    Load PowerPoint files
    WHY: python-pptx reads each slide's text
    We treat each slide as a separate document chunk
    """
    prs = Presentation(file_path)
    docs = []
    for slide_num, slide in enumerate(prs.slides):
        slide_text = f"Slide {slide_num + 1}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += shape.text + "\n"
        if slide_text.strip():
            docs.append(Document(
                page_content=slide_text,
                metadata={"source": file_path, "file_type": "PPTX",
                          "page": slide_num, "source_file": os.path.basename(file_path)}
            ))
    return docs

def load_any_file(file_path):
    """
    MASTER LOADER — detects file type and calls right loader
    WHY: Single entry point — app.py just calls this one function
    This is the Factory Pattern in software design
    """
    ext = os.path.splitext(file_path)[1].lower()
    print(f"📂 Loading {ext} file: {os.path.basename(file_path)}")

    if ext == ".pdf":
        return load_pdf(file_path)
    elif ext == ".docx":
        return load_docx(file_path)
    elif ext == ".txt":
        return load_txt(file_path)
    elif ext in [".csv", ".xlsx", ".xls"]:
        return load_csv_excel(file_path)
    elif ext == ".pptx":
        return load_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")